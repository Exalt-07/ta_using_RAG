# app/ingest.py
from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation
import hashlib
import shutil
from .retrieval import HybridRetriever # Relative import for your retriever class
from .multimodal import MultimodalEngine
class LectureProcessor:
    def __init__(self):
        self.retriever = HybridRetriever()
        self.vision_engine = MultimodalEngine()
    
    def process_lectures(self):
        # Check for model changes
        self._check_model_version()
        
        # Process text documents
        text_files = list(Path('data/raw').glob('*.*'))
        documents = self._extract_text(text_files)
        self.retriever.build_indices(documents)
        
        # Process images if available
        self._process_images()

    def _check_model_version(self):
        """Detect embedding model changes and reset ChromaDB if needed"""
        model_hash = hashlib.md5(
            self.retriever.text_model.encode("model_version_check").tobytes()
        ).hexdigest()
        
        hash_file = Path("data/processed/model_version.hash")
        
        if hash_file.exists():
            with open(hash_file, "r") as f:
                saved_hash = f.read()
            if saved_hash != model_hash:
                print("⚠️ Embedding model changed! Resetting vector store...")
                shutil.rmtree("data/processed/chromadb", ignore_errors=True)
        
        with open(hash_file, "w") as f:
            f.write(model_hash)

    def _process_images(self):
        image_dir = Path('data/raw/images')
        if image_dir.exists() and any(image_dir.iterdir()):
            print("Processing images...")
            self.vision_engine.process_images(image_dir)
        else:
            print("No images found in data/raw/images")

    def _extract_text(self, files):
        return [self._read_file(f) for f in files if f.is_file()]

    def _read_file(self, file_path):
        """Read text from supported file types"""
        try:
            if file_path.suffix.lower() == '.pdf':
                return self._read_pdf(file_path)
            elif file_path.suffix.lower() == '.pptx':
                return self._read_pptx(file_path)
            elif file_path.suffix.lower() == '.txt':
                return self._read_txt(file_path)
            else:
                print(f"⚠️ Unsupported file type: {file_path.suffix}")
                return ""
        except Exception as e:
            print(f"❌ Error reading {file_path.name}: {str(e)}")
            return ""

    def _read_pdf(self, file_path):
        reader = PdfReader(file_path)
        return " ".join(page.extract_text() or "" for page in reader.pages)

    def _read_pptx(self, file_path):
        prs = Presentation(file_path)
        return " ".join(
            shape.text for slide in prs.slides 
            for shape in slide.shapes if hasattr(shape, 'text')
        )

    def _read_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
